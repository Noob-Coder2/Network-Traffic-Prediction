import os
from io import BytesIO
from tkinter import messagebox

# Document merging libraries
from docx import Document
from pptx import Presentation
from pypdf import PdfWriter, PdfReader

# Optional import for Office COM integration
try:
    import comtypes.client
    OFFICE_INSTALLED = True
except (ImportError, OSError):
    OFFICE_INSTALLED = False

class FileProcessor:
    """
    Handles all file merging and processing logic. This runs in a separate
    thread to keep the UI responsive.
    """
    def __init__(self, app):
        self.app = app
        self.OFFICE_INSTALLED = OFFICE_INSTALLED

    def process_files(self, file_paths, output_filename, output_format):
        """
        The main entry point for the worker thread. Dispatches to the correct
        merge function based on the selected output format.
        """
        self.app.processing_queue.put(("max_progress", len(file_paths)))
        try:
            self.merge_documents(file_paths, output_filename, output_format)
        except Exception as e:
            self.app.processing_errors.append(f"A critical error occurred: {e}")
        finally:
            # Signal completion to the main thread
            self.app.processing_queue.put(("done", (file_paths, output_filename)))

    def merge_documents(self, file_paths, output_filename, output_format):
        """
        Selects and calls the appropriate merging function.
        """
        merge_function_map = {
            "TXT": self.merge_as_text,
            "PDF": self.merge_as_pdf,
            "DOCX": self.merge_as_docx,
            "PPTX": self.merge_as_pptx
        }
        merge_function = merge_function_map.get(output_format)

        if merge_function:
            merge_function(file_paths, output_filename)
        
        if self.app.cancel_requested:
            self.app.processing_queue.put(("status", "Operation cancelled."))
        else:
            status = f"Finished. Consolidated {len(file_paths) - len(self.app.processing_errors)} files."
            self.app.processing_queue.put(("status", status))

    def merge_as_text(self, file_paths, output_filename):
        """Merges a list of files into a single text file."""
        separator = "=" * 20
        with open(output_filename, 'w', encoding='utf-8') as outfile:
            for i, path in enumerate(file_paths):
                if self.app.cancel_requested: return
                self.app.processing_queue.put(("status", f"Processing: {os.path.basename(path)}"))
                self.app.processing_queue.put(("progress", i + 1))
                if self.app.include_headers.get():
                    outfile.write(f"\n{separator}\n# File: {path}\n{separator}\n\n")
                try:
                    with open(path, 'r', encoding='utf-8', errors='replace') as infile:
                        outfile.write(infile.read() + "\n")
                except Exception as e:
                    error_msg = f"Error reading file {path}: {e}"
                    outfile.write(f"# {error_msg}\n")
                    self.app.processing_errors.append(error_msg)

def merge_as_pdf(self, file_paths, output_filename):
    """Merges PDF, DOCX, and PPTX files into a single PDF."""
    if not self.OFFICE_INSTALLED and any(p.lower().endswith(('.docx', '.doc', '.pptx', '.ppt')) for p in file_paths):
        messagebox.showwarning("Office Not Found", "MS Office is required for converting Word/PowerPoint files to PDF. These files will be skipped.")

    writer = PdfWriter()  # Use PdfWriter
    temp_pdfs = []
    word, powerpoint = None, None
    try:
        for i, path in enumerate(file_paths):
            if self.app.cancel_requested: break
            self.app.processing_queue.put(("status", f"Converting: {os.path.basename(path)}"))
            self.app.processing_queue.put(("progress", i + 1))
            pdf_path = None

            if path.lower().endswith('.pdf'):
                pdf_path = path
            elif path.lower().endswith(('.docx', '.doc')) and self.OFFICE_INSTALLED:
                try:
                    if word is None: word = comtypes.client.CreateObject('Word.Application')
                    pdf_path = os.path.splitext(path)[0] + f"_{i}.pdf"
                    doc = word.Documents.Open(os.path.abspath(path))
                    doc.SaveAs(os.path.abspath(pdf_path), FileFormat=17)
                    doc.Close()
                    temp_pdfs.append(pdf_path)
                except Exception as e: self.app.processing_errors.append(f"Failed to convert DOCX {os.path.basename(path)}: {e}")
            elif path.lower().endswith(('.pptx', '.ppt')) and self.OFFICE_INSTALLED:
                try:
                    if powerpoint is None: powerpoint = comtypes.client.CreateObject('Powerpoint.Application')
                    pdf_path = os.path.splitext(path)[0] + f"_{i}.pdf"
                    ppt = powerpoint.Presentations.Open(os.path.abspath(path))
                    ppt.SaveAs(os.path.abspath(pdf_path), 32)  # 32 is PDF format
                    ppt.Close()
                    temp_pdfs.append(pdf_path)
                except Exception as e: self.app.processing_errors.append(f"Failed to convert PPTX {os.path.basename(path)}: {e}")
            else:
                self.app.processing_errors.append(f"Skipped unsupported file for PDF merge: {os.path.basename(path)}")
                continue

            if pdf_path:
                self.app.processing_queue.put(("status", f"Appending: {os.path.basename(path)}"))
                with open(pdf_path, 'rb') as f:
                    reader = PdfReader(f)
                    for page in reader.pages:
                        writer.add_page(page)  # Append each page

        if not self.app.cancel_requested:
            self.app.processing_queue.put(("status", "Saving final PDF..."))
            with open(output_filename, "wb") as f_out:
                writer.write(f_out)
    except Exception as e:
        self.app.processing_errors.append(f"A critical error occurred during PDF merge: {e}")
    finally:
        if word: word.Quit()
        if powerpoint: powerpoint.Quit()
        for temp_pdf in temp_pdfs:
            if os.path.exists(temp_pdf):
                try:
                    os.remove(temp_pdf)
                except OSError as e:
                    self.app.processing_errors.append(f"Could not remove temp file {temp_pdf}: {e}")


    def merge_as_docx(self, file_paths, output_filename):
        """Merges multiple DOCX files into one."""
        master_doc = Document()
        if file_paths and self.app.include_headers.get():
            master_doc.add_heading('Consolidated Document', 0)
        
        for i, path in enumerate(file_paths):
            if self.app.cancel_requested: return
            self.app.processing_queue.put(("status", f"Processing: {os.path.basename(path)}"))
            self.app.processing_queue.put(("progress", i + 1))
            
            if path.lower().endswith('.docx'):
                try:
                    if i > 0 or not self.app.include_headers.get():
                        master_doc.add_page_break()
                    if self.app.include_headers.get():
                        master_doc.add_heading(f"Content from: {os.path.basename(path)}", level=1)
                    
                    sub_doc = Document(path)
                    for element in sub_doc.element.body:
                        master_doc.element.body.append(element)
                except Exception as e:
                    self.app.processing_errors.append(f"Could not process DOCX file {os.path.basename(path)}: {e}")
            else:
                self.app.processing_errors.append(f"Skipped non-DOCX file: {os.path.basename(path)}")
        
        if not self.app.cancel_requested:
            master_doc.save(output_filename)

    def merge_as_pptx(self, file_paths, output_filename):
        """Merges multiple PPTX files into one."""
        master_ppt = Presentation()
        for i, path in enumerate(file_paths):
            if self.app.cancel_requested: return
            self.app.processing_queue.put(("status", f"Processing: {os.path.basename(path)}"))
            self.app.processing_queue.put(("progress", i + 1))
            
            if path.lower().endswith('.pptx'):
                try:
                    sub_ppt = Presentation(path)
                    for slide in sub_ppt.slides:
                        # Add a new slide with the same layout
                        new_slide = master_ppt.slides.add_slide(slide.slide_layout)
                        # Copy shapes from old slide to new slide
                        for shape in slide.shapes:
                            if shape.has_text_frame:
                                new_shape = new_slide.shapes.add_textbox(shape.left, shape.top, shape.width, shape.height)
                                new_shape.text_frame.text = shape.text_frame.text
                            elif shape.shape_type == 13: # Picture
                                try:
                                    image_stream = BytesIO(shape.image.blob)
                                    new_slide.shapes.add_picture(image_stream, shape.left, shape.top, width=shape.width, height=shape.height)
                                except Exception as img_e: self.app.processing_errors.append(f"Could not copy image from {os.path.basename(path)}: {img_e}")
                            # Note: This is a simplified copy. More complex shapes, tables, charts etc. are not copied.
                except Exception as e:
                    self.app.processing_errors.append(f"Could not process {os.path.basename(path)}: {e}")
            else:
                self.app.processing_errors.append(f"Skipped non-PPTX file: {os.path.basename(path)}")
        
        if not self.app.cancel_requested:
            master_ppt.save(output_filename)
